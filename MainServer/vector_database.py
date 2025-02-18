from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv
from pptx import Presentation
from flask import jsonify
from docx import Document
import numpy as np
import PyPDF2
import pickle
import faiss
import os

load_dotenv()
DOC_PATH = os.getenv("DOC_PATH")
V_DB_PATH = os.getenv("V_DB_PATH")

def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """
    Breaks a string into overlapping text chunks.
   
    Returns:
        list: A list of text chunk strings.
    """
    words = text.split()
    chunks = []
    index = 0
    while index < len(words):
        chunk_words = words[index: index + chunk_size]
        chunks.append(" ".join(chunk_words))
        if index + chunk_size >= len(words):
            break
        index += chunk_size - chunk_overlap
    return chunks

def generate_chunks_from_file(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """   
    Returns:
        List[dict]: List of dictionaries containing:
            - 'chunk': Text chunk.
            - 'page': Page/Slide/chunk sequence number.
    """
    chunks = []
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, start=1):
                    page_text = page.extract_text() or ""
                    if page_text:
                        page_chunks = chunk_text(page_text, chunk_size, chunk_overlap)
                        for chunk in page_chunks:
                            chunks.append({"chunk": chunk, "page": page_num})
        elif ext == ".docx":
            doc = Document(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            if full_text:
                doc_chunks = chunk_text(full_text, chunk_size, chunk_overlap)
                for i, chunk in enumerate(doc_chunks, start=1):
                    # For files without inherent pages, using sequential chunk numbers.
                    chunks.append({"chunk": chunk, "page": i})
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                slide_text = "\n".join(slide_text)
                if slide_text:
                    slide_chunks = chunk_text(slide_text, chunk_size, chunk_overlap)
                    for chunk in slide_chunks:
                        chunks.append({"chunk": chunk, "page": slide_num})
        else:
            # For txt, md, etc. files
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                if text:
                    text_chunks = chunk_text(text, chunk_size, chunk_overlap)
                    for i, chunk in enumerate(text_chunks, start=1):
                        chunks.append({"chunk": chunk, "page": i})
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
    
    return chunks

def process_multiple_files(doc_dir: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """    
    Returns:
        List[dict]: A list of dictionaries. Each dictionary contains:
            - 'file': File name.
            - 'chunk': The text chunk.
            - 'page': Page/Slide/chunk sequence number.
    """
    chunks_data = []
    
    for file_name in os.listdir(doc_dir):
        file_path = os.path.join(doc_dir, file_name)
        chunks = generate_chunks_from_file(file_path, chunk_size, chunk_overlap)
        for chunk_data in chunks:
            # Here the key 'file' stores the filename.
            chunks_data.append({
                "file": file_name,
                "chunk": chunk_data["chunk"],
                "page": chunk_data["page"]
            })

    return chunks_data

def embed_chunks(chunks_data: list, model_name: str = "all-MiniLM-L6-v2", convert_to_tensor: bool = True) -> list:
    """
    Returns:
        List[dict]: The input list with an added key 'embedding' for each dictionary.
    """
    # Load the pre-trained model
    model = SentenceTransformer(model_name)
    
    # Extract all text chunks for encoding
    texts = [item["chunk"] for item in chunks_data]
    
    # Generate embeddings for all chunks at once
    embeddings = model.encode(texts, convert_to_tensor=convert_to_tensor)
    
    # Attach each embedding to its corresponding dictionary
    for i, item in enumerate(chunks_data):
        item["embedding"] = embeddings[i]
    
    return chunks_data

def save_chunks_to_faiss(chunks_data: list,
                         index_file: str = "faiss_index.idx",
                         metadata_file: str = "metadata.pkl") -> None:
    """
    Save embeddings from chunks_data into a FAISS vector index, creating necessary files if they don't exist.
    """

    index_file = os.path.join(os.getcwd(), index_file)
    metadata_file = os.path.join(os.getcwd(), metadata_file)

    os.makedirs(os.path.dirname(index_file), exist_ok=True)  # Ensure directory exists

    embeddings = []
    metadata = []
    ids = []

    # Check if FAISS index and metadata file exist
    if os.path.exists(index_file):
        index = faiss.read_index(index_file)
    else:
        index = None  # Will create a new one

    if os.path.exists(metadata_file):
        with open(metadata_file, "rb") as f:
            metadata = pickle.load(f)
    else:
        metadata = []

    # Determine the starting ID
    current_max_id = max((item["id"] for item in metadata), default=-1)

    for i, item in enumerate(chunks_data):
        try:
            # Convert embedding to NumPy array if it's a PyTorch tensor
            emb = item["embedding"].cpu().numpy() if hasattr(item["embedding"], "cpu") else item["embedding"]
            embeddings.append(emb)

            # Assign a new unique ID
            new_id = current_max_id + i + 1
            ids.append(new_id)

            # Save metadata
            metadata.append({
                "file": item["file"],
                "chunk": item["chunk"],
                "page": item.get("page", "Unknown"),
                "id": new_id
            })
        except Exception as e:
            print(f"Error processing chunk {i}: {e}")

    if not embeddings:
        print("No embeddings found. Nothing to save.")
        return

    # Stack embeddings into a 2D NumPy array (float32)
    embeddings_matrix = np.vstack(embeddings).astype("float32")
    
    # Normalize embeddings for cosine similarity
    faiss.normalize_L2(embeddings_matrix)

    embedding_dim = embeddings_matrix.shape[1]
    ids_np = np.array(ids, dtype=np.int64)

    # Create index if it doesn't exist
    if index is None:
        index_flat = faiss.IndexFlatIP(embedding_dim)
        index = faiss.IndexIDMap(index_flat)

    # Add new embeddings
    index.add_with_ids(embeddings_matrix, ids_np)

    # Save updated FAISS index and metadata
    faiss.write_index(index, index_file)
    with open(metadata_file, "wb") as f:
        pickle.dump(metadata, f)

    print(f"Saved FAISS index with {index.ntotal} vectors to '{index_file}'.\n")

def delete_chunks_from_file(
                            file_name: str,
                            has_multiple_files: bool = True,
                            index_file: str = "faiss_index.idx",
                            metadata_file: str = "metadata.pkl") -> None:
    """
    Delete all vector chunks corresponding to a given file from the FAISS vector database.
    """

    index_file = f"{os.getcwd()}\\{index_file}"
    metadata_file = f"{os.getcwd()}\\{metadata_file}"

    # Load the FAISS index
    index = faiss.read_index(index_file)
    
    # Load the existing metadata
    try:
        with open(metadata_file, "rb") as f:
            metadata = pickle.load(f)
    except FileNotFoundError:
        print(f"Metadata file '{metadata_file}' not found.")
        return
    
    ids_to_delete = []
    updated_metadata = []
    
    # Identify vectors (by id) that belong to the specified file.
    for item in metadata:
        if item.get("file") == file_name:
            if "id" in item:
                ids_to_delete.append(item["id"])
            else:
                print("Error: Metadata item missing 'id'. Cannot delete without vector IDs.")
                return
        else:
            updated_metadata.append(item)
    
    if not ids_to_delete:
        return jsonify({"error": f"No chunks found for file {file_name}"}), 409
    
    # Convert the list of ids to a NumPy array of type int64
    ids_to_delete = np.array(ids_to_delete, dtype=np.int64)
    
    # Remove the vectors corresponding to these ids from the FAISS index.
    index.remove_ids(ids_to_delete)
    
    # Save the updated index and metadata back to disk
    faiss.write_index(index, index_file)
    with open(metadata_file, "wb") as f:
        pickle.dump(updated_metadata, f)
    
    print(f"Deleted {len(ids_to_delete)} chunks from file '{file_name}'.")
    if (not has_multiple_files):    # if it was the only file to be deleted, remove the folder
        print(f"deleted {index_file}")
        print(f"deleted {metadata_file}")
        
        os.remove(index_file)
        os.remove(metadata_file)

if __name__ == "__main__":
    course_codes = ['CS307']
    for course_code in course_codes:
        doc_directory = f"{DOC_PATH}\\{course_code}" 
        # Process PDFs and generate text chunks
        chunks_data = process_multiple_files(doc_directory, chunk_size=1000, chunk_overlap=200)
        print(f"Extracted {len(chunks_data)} chunks for {course_code}.")
    
        # Generate embeddings for each chunk using Sentence Transformers
        chunks_data = embed_chunks(chunks_data, model_name="all-MiniLM-L6-v2", convert_to_tensor=True)

        # Save the chunks and their embeddings to a FAISS vector database (with metadata)
        save_chunks_to_faiss(chunks_data, index_file= f"{V_DB_PATH}\\{course_code}_faiss_index.idx", metadata_file=f"{V_DB_PATH}\\{course_code}_metadata.pkl")
        print(f"Processed the documents for {course_code}.")